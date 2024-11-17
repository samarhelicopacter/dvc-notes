import requests
import streamlit as st
from typing import Dict, Optional, List
from datetime import datetime
import json
from PIL import Image
import io
import base64
from pathlib import Path
import pytesseract
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import tempfile
import os

# Set page config must be the first Streamlit command
st.set_page_config(
    page_title="DVC NoteSmart",
    page_icon="📝",
    layout="wide",
    initial_sidebar_state="expanded"
)

class DVCNoteSmart:
    def __init__(self, api_key):
        self.api_key = api_key
        self.base_url = "https://api.anthropic.com/v1/messages"
        self.headers = {
            "anthropic-version": "2023-06-01",
            "content-type": "application/json",
            "x-api-key": self.api_key
        }
        # Enhanced definitions dictionary with more CCC-specific terms
        self.ccc_definitions = {
            "Vision 2030": """The California Community Colleges Vision 2030 initiative - A strategic plan focusing on:
            - Closing equity gaps
            - Improving transfer rates
            - Increasing workforce alignment
            - Enhancing student success metrics
            - Implementing innovative teaching practices""",
            
            "AB928": """Assembly Bill 928 (Student Transfer Achievement Reform Act of 2021):
            - Establishes a singular lower-division general education pathway
            - Creates the California General Education Transfer Curriculum (CalGETC)
            - Requires implementation by Fall 2025
            - Aims to streamline transfer between CCC, CSU, and UC systems""",
            
            "FTES": """Full-Time Equivalent Student:
            - Key funding metric for California Community Colleges
            - Calculated as 525 hours of student instruction annually
            - Used for state funding allocations
            - Critical for budget planning and resource allocation""",
            
            "CalGETC": """California General Education Transfer Curriculum:
            - New unified GE pattern for UC/CSU transfer starting Fall 2025
            - Replaces IGETC and CSU GE Breadth
            - Requires 34 semester units
            - Includes Ethnic Studies requirement
            - Designed to simplify transfer process""",
            
            "IGETC": """Intersegmental General Education Transfer Curriculum:
            - Current transfer pattern for UC/CSU
            - Being replaced by CalGETC in Fall 2025
            - Important for transfer planning and articulation""",
            
            "Title 5": """California Code of Regulations, Title 5:
            - Governs California Community Colleges
            - Establishes educational standards
            - Defines degree and certificate requirements
            - Sets policies for curriculum and instruction"""
        }

    def process_file(self, uploaded_file) -> str:
        """Process a single file and return its content as text."""
        try:
            # Get file extension
            file_extension = os.path.splitext(uploaded_file.name)[1].lower()
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                # Write uploaded file content to temporary file
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name
                
            # Extract content based on file type
            content = ""
            
            # Handle PDFs
            if file_extension == '.pdf':
                with open(tmp_file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
            
            # Handle Word documents
            elif file_extension in ['.docx', '.doc']:
                doc = Document(tmp_file_path)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Handle PowerPoint
            elif file_extension in ['.pptx', '.ppt']:
                prs = Presentation(tmp_file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            
            # Handle Excel
            elif file_extension in ['.xlsx', '.xls']:
                df = pd.read_excel(tmp_file_path)
                content = df.to_string()
            
            # Handle Images - Updated this part
            elif file_extension in ['.png', '.jpg', '.jpeg']:
                try:
                    # Open and convert image to RGB
                    with Image.open(tmp_file_path) as image:
                        # Convert to RGB if needed
                        if image.mode != 'RGB':
                            image = image.convert('RGB')
                        # Save as temporary PNG
                        temp_png = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                        image.save(temp_png.name)
                        # Process with tesseract
                        content = pytesseract.image_to_string(Image.open(temp_png.name))
                        # Clean up temporary PNG
                        os.unlink(temp_png.name)
                except Exception as img_error:
                    st.warning(f"Image processing error: {str(img_error)}")
                    content = "Error: Could not extract text from image."
            
            # Clean up temporary file
            os.unlink(tmp_file_path)
            
            return content
        
        except Exception as e:
            st.warning(f"Error processing {uploaded_file.name}: {str(e)}")
            return ""

    def process_all_files(self, uploaded_files) -> str:
        """Process all uploaded files and combine their content."""
        all_content = ""
        
        for uploaded_file in uploaded_files:
            content = self.process_file(uploaded_file)
            if content:
                all_content += f"\n\n=== Content from {uploaded_file.name} ===\n"
                all_content += content
        
        return all_content
    
    def setup_page_style(self):
        """Setup custom CSS styling for the application."""
        st.markdown("""
            <style>
            /* DVC brand colors */
            :root {
                --dvc-green: #006341;
                --dvc-gold: #9B8B5B;
                --dvc-gray: #58595B;
                --dvc-light-gray: #f0f2f6;
            }
        
            .main-header {
                background: linear-gradient(135deg, var(--dvc-green), var(--dvc-gold));
                padding: 2rem;
                border-radius: 10px;
                color: white;
                margin-bottom: 2rem;
                text-align: center;
                box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
            }
            
            .sub-header {
                color: var(--dvc-green);
                margin: 1.5rem 0;
                font-weight: 600;
                font-size: 1.25rem;
            }
            
            .stButton > button {
                background-color: var(--dvc-green) !important;
                color: white !important;
                font-weight: 600 !important;
                padding: 0.5rem 1rem !important;
                border-radius: 5px !important;
                border: none !important;
                transition: all 0.3s ease !important;
            }
            
            .stButton > button:hover {
                background-color: var(--dvc-gold) !important;
                transform: translateY(-2px);
            }
            
            .info-box {
                background-color: var(--dvc-light-gray);
                padding: 1rem;
                border-radius: 5px;
                margin: 1rem 0;
                border-left: 4px solid var(--dvc-green);
            }
            
            .file-preview {
                border: 1px solid var(--dvc-gray);
                padding: 1rem;
                border-radius: 5px;
                margin: 0.5rem 0;
            }
            
            .stFileUploader {
                border: 2px dashed var(--dvc-gray);
                padding: 1rem;
                border-radius: 5px;
            }
            
            .stTextArea textarea {
                border: 1px solid var(--dvc-gray);
                border-radius: 5px;
                padding: 0.5rem;
            }
            
            .stSelectbox select {
                border-radius: 5px;
                border-color: var(--dvc-gray);
            }
            
            .stRadio > div {
                flex-direction: row !important;
                gap: 2rem !important;
            }
            
            .stRadio label {
                background-color: var(--dvc-light-gray);
                padding: 0.5rem 1rem;
                border-radius: 5px;
                cursor: pointer;
            }
            
            .stRadio label:hover {
                background-color: var(--dvc-gold);
                color: white;
            }
            </style>
        """, unsafe_allow_html=True)

    def _create_meeting_prompt(self, text: str, meeting_type: str, context: Dict, detail_level: str) -> str:
        """Create a detailed prompt for the meeting notes."""
        detail_instructions = {
            "standard": """Create clear, concise notes that capture the main points and decisions.""",
            "detailed": """Create comprehensive notes with:
            - Extended discussion summaries
            - Detailed background context
            - Supporting data and research
            - Complex policy implications
            - Historical context when relevant
            - Alternative viewpoints discussed
            - Stakeholder concerns and responses
            - Related initiatives and connections
            - Potential long-term impacts"""
        }
        
        base_prompt = f"""You are a professional note-taker for {meeting_type} meetings at Diablo Valley College. 
        {detail_instructions[detail_level]}
        Create detailed, professional notes from this transcript following these requirements:

        Meeting Context:
        - Date: {context.get('date', 'Not specified')}
        - Type: {meeting_type}
        
        Referenced Materials:
        {self._format_materials_list(context.get('uploaded_materials', []))}

        Required Sections:

        1. Meeting Overview
        - Date
        - Purpose and context
        - Major takeaways
        - Strategic alignment with Vision 2030 and DVC goals

        2. Discussion Points
        - Clear, bolded headers for each topic
        - Detailed summaries with metrics and policies
        - Definitions of technical terms, laws, or systems
        - Connection to Vision 2030 and DVC priorities
        - Visual data interpretation when provided

        3. Strategic Alignment
        - Vision 2030 goals (equity gaps, transfer rates, workforce)
        - DVC's mission and strategic objectives
        - Specific initiatives and their alignment

        4. Action Items
        - Responsible parties
        - Clear deadlines
        - Required resources
        - Success metrics
        - Follow-up plans

        5. Next Steps
        - Upcoming deadlines
        - Future meetings
        - Required preparation

        Formatting Requirements:
        - Use markdown formatting
        - Bold (**) for headers and key terms
        - Use bullet points for lists
        - Include clear section breaks
        - Maintain professional tone
        - Format dates consistently
        - Number action items
        - Use tables for complex data

        Key Terms to Reference:
        {json.dumps(self.ccc_definitions, indent=2)}

        Meeting transcript to analyze:

        {text}"""
        print (base_prompt)
        return base_prompt

    def _format_materials_list(self, materials: List[str]) -> str:
        """Format the list of uploaded materials."""
        if not materials:
            return "No additional materials provided"
        
        formatted_list = "Uploaded Materials:\n"
        for idx, material in enumerate(materials, 1):
            formatted_list += f"{idx}. {material}\n"
        return formatted_list
    
    def process_notes(self, text: str, meeting_type: str = "General", context: Dict = None, detail_level: str = "standard") -> Optional[Dict]:
        """Process text using Claude API to generate professional meeting notes."""
        try:
            if context is None:
                context = {}
                
            # ADD THIS NEW SECTION HERE - This is the only new part we're adding
            uploaded_files = context.get('uploaded_materials', [])
            if uploaded_files:
                file_content = self.process_all_files(uploaded_files)
                if file_content:
                    text += "\n\nContent from uploaded materials:\n" + file_content
            # END OF NEW SECTION

            system_prompt = """You are a professional note-taker for California Community Colleges, specializing in creating 
            comprehensive, detailed meeting notes with rich narrative descriptions. Focus on:
            1. Providing extensive context for each discussion point
            2. Capturing the flow and evolution of discussions
            3. Explaining rationale behind decisions
            4. Using roles/functions instead of names
            5. Including specific examples while maintaining anonymity
            6. Preserving technical details and data points
            
            Format notes with:
            - Clear section headings
            - Detailed bullet points with complete context
            - Rich narrative descriptions
            - Specific examples and data
            - Decision rationale
            - Next steps and implications"""

            # Enhanced meeting prompt
            meeting_prompt = f"""Create detailed, narrative-style notes from this transcript following these requirements:

            Meeting Context:
            - Type: {meeting_type}
            - Referenced Materials: {self._format_materials_list(context.get('uploaded_materials', []))}

            Required Sections:

            1. Quick Navigation
            - Links to major topics
            - Key decisions summary
            - Critical deadlines

            2. Discussion Points
            For each major topic, include:
            • **Context & Background:** Detailed explanation of why this topic was discussed
            • **Key Updates:** Rich narrative descriptions of main points
            • **Challenges & Solutions:** Thorough exploration of issues raised and solutions proposed
            • **Decisions & Rationale:** Complete context for why decisions were made
            • **Implementation Details:** Specific steps and considerations
            • **Examples & Data:** Relevant numbers, scenarios, or cases discussed
            
            3. Action Items
            By responsible area (not individual):
            • Specific tasks with complete context
            • Dependencies and requirements
            • Resource needs
            • Success criteria

            4. Next Steps
            • Upcoming work required
            • Preparation needed
            • Dependencies and timelines

            Format Requirements:
            - Use rich narrative bullet points that tell complete stories
            - Include specific examples and data points
            - Maintain anonymity while preserving context
            - Bold key terms and concepts
            - Preserve technical details and numbers
            - Group related items logically

            Reference these terms and definitions when relevant:
            {json.dumps(self.ccc_definitions, indent=2)}

            Meeting transcript to analyze:

            {text}"""
            print (system_prompt)
            print (meeting_prompt)
            data = {
                "model": "claude-3-sonnet-20240229",
                "system": system_prompt,
                "messages": [
                    {
                        "role": "user",
                        "content": meeting_prompt
                    }
                ],
                "max_tokens": 4096,
                "temperature": 0.7
            }

            try:
                response = requests.post(
                    self.base_url,
                    headers=self.headers,
                    json=data,
                    timeout=60
                )
                response.raise_for_status()
                
                content = response.json()['content'][0]['text']
                
                return {
                    "summary": content,
                    "meeting_type": meeting_type,
                    "context": context,
                    "detail_level": detail_level,
                    "model_version": "claude-3-sonnet-20240229"
                }
                
            except requests.exceptions.Timeout:
                st.error("Request timed out. Please try again.")
                return None
            except requests.exceptions.RequestException as e:
                st.error(f"API Error: {str(e)}")
                self._handle_api_error(response)
                return None
            
        except Exception as e:
            st.error(f"Error processing notes: {str(e)}")
            return None

    def _handle_api_error(self, response):
        """Handle API errors with detailed feedback."""
        if hasattr(response, 'status_code'):
            if response.status_code == 429:
                st.warning("Rate limit exceeded. Please wait a moment before trying again.")
            elif response.status_code == 401:
                st.error("Invalid API key. Please check your credentials.")
            try:
                error_detail = response.json()
                st.error(f"API Error Details: {error_detail}")
            except:
                pass

def main():
    # Initialize API key
    api_key = st.secrets["ANTHROPIC_API_KEY"]
    processor = DVCNoteSmart(api_key)

    # Setup page styling
    processor.setup_page_style()

    # Main header
    st.markdown("""
        <div class="main-header">
            <h1>DVC NoteSmart</h1>
            <p>Smart Meeting Notes Assistant for DVC</p>
        </div>
    """, unsafe_allow_html=True)

    # Sidebar content
    with st.sidebar:
        st.header("Settings")
        meeting_type = st.selectbox(
            "Meeting Type",
            ["Academic Senate", "Department Meeting", "Committee Meeting",
            "Professional Development", "Other"]
        )
        
        st.markdown("### About DVC NoteSmart")
        st.markdown("""
            Smart meeting notes assistant that aligns with DVC's mission, values, and strategic goals.
        """)
        
        with st.expander("DVC Mission"):
            st.markdown("""
                **Our Mission**
                
                We inspire, educate, and empower students to transform their lives and their communities.
                We guide students to achieve their goals by:
                - Awarding degrees and certificates
                - Preparing for transfer to 4-year institutions
                - Facilitating career advancement
                - Fostering personal growth
            """)
        
        with st.expander("Institutional Learning Outcomes"):
            st.markdown("""
                **Core Competencies**
                - Communication and Collaboration
                - Empathy Mindset
                - Growth Mindset
                - Information and Technology Fluency
                - Solution Mindset
            """)
        
        with st.expander("Strategic Frameworks"):
            st.markdown("""
                **Vision 2030**
                A collaborative action plan providing focus, equity, and direction to our community colleges.
                
                **4CD District Strategic Plan**
                Aligned with the District Strategic Plan 2020-2025 goals and objectives.
                
                **General Education Outcomes**
                Supporting DVC's commitment to:
                - Effective communication
                - Critical analysis and problem-solving
                - Cultural and artistic appreciation
                - Scientific inquiry and understanding
                - Health and wellness promotion
            """)
        
        # File upload section
        st.header("Meeting Materials")
        uploaded_files = st.file_uploader(
            "Upload Meeting Materials",
            type=['pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'png', 'jpg', 'jpeg'],
            accept_multiple_files=True,
            help="Upload any meeting materials (documents, presentations, images, etc.)"
    )

    # Main content area
    st.markdown('<p class="sub-header">Meeting Details</p>', unsafe_allow_html=True)
    
    # Meeting metadata
    meeting_date = st.date_input(
        "Meeting Date",
        format="MM/DD/YYYY"
    )
    
    # Add detail level selector
    detail_level = st.radio(
        "Note Detail Level",
        ["standard", "detailed"],
        format_func=lambda x: "Standard Notes" if x == "standard" else "Detailed Notes (Extended Analysis)",
        help="Choose the level of detail for your meeting notes. Detailed notes include extended discussion summaries, policy implications, and historical context."
    )

    # Text input for transcript
    st.markdown('<p class="sub-header">Meeting Transcript</p>', unsafe_allow_html=True)
    text_input = st.text_area(
        "Paste your meeting transcript below:",
        height=300,
        help="Include as much detail as possible for best results"
    )

    # Process button
    if st.button("Generate Notes", use_container_width=True):
        if text_input and text_input.strip():
            with st.spinner("🔄 Processing your notes... (this may take a few minutes)"):
                context = {
                    "date": meeting_date.strftime("%m/%d/%Y"),
                    "uploaded_materials": uploaded_files if uploaded_files else []
                }
                
                notes_data = processor.process_notes(text_input, meeting_type, context, detail_level)
                if notes_data:
                    st.markdown('<p class="sub-header">Generated Meeting Notes</p>', 
                            unsafe_allow_html=True)
                    st.markdown(notes_data["summary"])

if __name__ == "__main__":
    main()